import os
import argparse
from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def html_to_ppt(html_folder, output_ppt):
    global prs
    # 创建演示文稿
    prs = Presentation()
    
    # 获取所有HTML文件并按页码排序
    html_files = [f for f in os.listdir(html_folder) if f.endswith('.html')]
    # 按页码排序（假设文件名格式为page_X.html）
    html_files.sort(key=lambda x: int(x.split('_')[1].split('.')[0]))
    
    for html_file in html_files:
        html_path = os.path.join(html_folder, html_file)
        
        # 读取HTML文件
        with open(html_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # 解析HTML
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # 创建幻灯片
        slide_layout = prs.slide_layouts[1]  # 使用带标题和内容的布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题（从HTML的h1标签获取）
        title = slide.shapes.title
        h1_tag = soup.find('h1')
        if h1_tag:
            title.text = h1_tag.get_text(strip=True)
        else:
            title.text = f"{html_file.replace('.html', '')}"
        
        # 获取内容占位符
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.word_wrap = True
        
        # 提取HTML中的文本内容和格式
        body_content = soup.body
        if body_content:
            # 递归处理HTML元素
            process_element(body_content, tf)


def process_element(element, text_frame):
    # 处理文本节点
    if isinstance(element, str):
        text = element.strip()
        if text and text_frame.paragraphs:
            # 添加到最后一个段落
            text_frame.paragraphs[-1].text += text
        return
    
    # 处理标签元素
    if isinstance(element, Tag):
        # 对于标题和段落标签，创建新段落
        if element.name in ['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            ppt_p = text_frame.add_paragraph()
            
            # 设置段落格式
            if element.name.startswith('h'):
                heading_level = int(element.name[1])
                if heading_level == 1:
                    ppt_p.font.size = Pt(24)
                elif heading_level == 2:
                    ppt_p.font.size = Pt(20)
                elif heading_level == 3:
                    ppt_p.font.size = Pt(18)
                else:
                    ppt_p.font.size = Pt(16)
                ppt_p.font.bold = True
            else:
                ppt_p.font.size = Pt(14)
            
            # 处理子元素
            for child in element.children:
                process_child(child, ppt_p)
        
        # 处理列表
        elif element.name in ['ul', 'ol']:
            for li in element.find_all('li', recursive=False):
                ppt_p = text_frame.add_paragraph()
                ppt_p.font.size = Pt(14)
                
                # 设置列表符号
                if element.name == 'ul':
                    ppt_p.level = 1
                else:
                    ppt_p.level = 1
                    # 数字列表在pptx中会自动生成
                
                # 处理列表项内容
                for child in li.children:
                    process_child(child, ppt_p)
        
        # 处理其他标签，递归处理子元素
        else:
            for child in element.children:
                process_element(child, text_frame)


def process_child(child, paragraph):
    # 处理文本节点
    if isinstance(child, str):
        text = child.strip()
        if text:
            paragraph.text += text
    
    # 处理标签元素
    elif isinstance(child, Tag):
        # 处理加粗
        if child.name == 'strong' or child.name == 'b':
            run = paragraph.add_run()
            run.text = child.get_text(strip=True)
            run.font.bold = True
        
        # 处理斜体
        elif child.name == 'em' or child.name == 'i':
            run = paragraph.add_run()
            run.text = child.get_text(strip=True)
            run.font.italic = True
        
        # 处理下划线
        elif child.name == 'u':
            run = paragraph.add_run()
            run.text = child.get_text(strip=True)
            run.font.underline = True
        
        # 处理其他标签
        else:
            paragraph.text += child.get_text(strip=True)
    



if __name__ == '__main__':
    # 创建命令行参数解析器
    parser = argparse.ArgumentParser(description='将HTML文件转换为PPT')
    parser.add_argument('html_folder', type=str, help='包含HTML文件的文件夹路径')
    parser.add_argument('output_ppt', type=str, help='输出PPT文件路径（包含文件名）')
    
    # 解析参数
    args = parser.parse_args()
    
    # 调用函数
    html_to_ppt(args.html_folder, args.output_ppt)
    
    # 保存PPT
    prs.save(args.output_ppt)
    print(f"PPT已保存至: {args.output_ppt}")